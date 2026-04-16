import json
import os
import re
from typing import List

import pandas as pd
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from pptx import Presentation

from ingestion.element_chunker import ElementChunker
from pipeline_logger import get_logger
from processing.chunker import TextChunker
from processing.cleaner import TextCleaner
from processing.language_detector import LanguageDetector


logger = get_logger("structured_loader")


class StructuredFileParser:
    def __init__(self):
        self.cleaner = TextCleaner()
        self.lang_detector = LanguageDetector()
        self.chunker = TextChunker()
        self.element_chunker = ElementChunker()

    def parse_text_like(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)
        extension = os.path.splitext(file_path)[1].lower()

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            text = f.read()

        return self._to_chunks(text=text, source=file_name, page=1, extension=extension)

    def parse_docx(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)
        document = DocxDocument(file_path)

        paragraphs = [p.text.strip() for p in document.paragraphs if (p.text or "").strip()]
        markdown_like = []
        for paragraph in document.paragraphs:
            raw = (paragraph.text or "").strip()
            if not raw:
                continue
            style_name = str(getattr(paragraph.style, "name", "") or "")
            if style_name.lower().startswith("heading"):
                markdown_like.append(f"## {raw}")
            else:
                markdown_like.append(raw)

        text = "\n\n".join(markdown_like or paragraphs)
        return self._to_chunks(text=text, source=file_name, page=1, extension=".md")

    def parse_html(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            raw = f.read()

        soup = BeautifulSoup(raw, "html.parser")
        text = soup.get_text(separator="\n")
        return self._to_chunks(text=text, source=file_name, page=1, extension=".html")

    def parse_json(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            obj = json.load(f)

        lines = []

        def walk(node, prefix=""):
            if isinstance(node, dict):
                for key, value in node.items():
                    next_prefix = f"{prefix}.{key}" if prefix else str(key)
                    walk(value, next_prefix)
            elif isinstance(node, list):
                for idx, value in enumerate(node):
                    next_prefix = f"{prefix}[{idx}]"
                    walk(value, next_prefix)
            else:
                lines.append(f"{prefix}: {node}")

        walk(obj)
        text = "\n".join(lines)
        return self._to_chunks(text=text, source=file_name, page=1, extension=".json")

    def parse_csv(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)
        df = pd.read_csv(file_path, dtype=str, keep_default_na=False)

        if df.empty:
            return []

        sample_text = " ".join(df.astype(str).head(15).to_string(index=False).split())
        language = self.lang_detector.detect_language(sample_text) if sample_text else "en"

        return self.element_chunker.chunk_table_dataframe(
            df=df,
            source=file_name,
            page=1,
            language=language,
            table_id="csv_table_1",
            section_hint="csv_data",
            rows_per_chunk=20,
        )

    def parse_excel(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)

        chunks = []
        sheet_map = pd.read_excel(file_path, sheet_name=None, dtype=str)

        for sheet_idx, (sheet_name, df) in enumerate(sheet_map.items(), start=1):
            if df is None or df.empty:
                continue

            df = df.fillna("")
            sample_text = " ".join(df.astype(str).head(12).to_string(index=False).split())
            language = self.lang_detector.detect_language(sample_text) if sample_text else "en"
            sheet_chunks = self.element_chunker.chunk_table_dataframe(
                df=df,
                source=file_name,
                page=sheet_idx,
                language=language,
                table_id=f"excel_{sheet_name}",
                section_hint=f"sheet:{sheet_name}",
                rows_per_chunk=20,
            )
            chunks.extend(sheet_chunks)

        return chunks

    def parse_pptx(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)
        presentation = Presentation(file_path)

        chunks = []

        for slide_idx, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and (shape.text or "").strip():
                    texts.append(shape.text.strip())

            slide_text = "\n".join(texts)
            slide_chunks = self._to_chunks(text=slide_text, source=file_name, page=slide_idx, extension=".md")
            chunks.extend(slide_chunks)

        return chunks

    def parse_ppt(self, file_path: str, source_name: str = None) -> List[dict]:
        file_name = source_name or os.path.basename(file_path)

        with open(file_path, "rb") as f:
            data = f.read()

        decoded = data.decode("latin-1", errors="ignore")
        text_bits = re.findall(r"[A-Za-z0-9\u0900-\u097F\u0980-\u09FF\u0B80-\u0BFF\u0C00-\u0C7F\u0D00-\u0D7F\s\.,:;\-_'\"\(\)\/]{6,}", decoded)
        text = "\n".join(text_bits)

        logger.warning("Legacy .ppt parsing is best-effort text extraction only | source=%s", file_name)
        return self._to_chunks(text=text, source=file_name, page=1, extension=".txt")

    def _to_chunks(self, text: str, source: str, page: int, extension: str = "") -> List[dict]:
        cleaned = self.cleaner.clean(text or "")
        if not cleaned.strip():
            return []

        language = self.lang_detector.detect_language(cleaned)

        ext = (extension or "").lower()
        if ext in {".md", ".markdown"}:
            return self.element_chunker.chunk_markdown(
                text=cleaned,
                source=source,
                page=page,
                language=language,
            )

        return self.element_chunker.chunk_generic(
            text=cleaned,
            source=source,
            page=page,
            language=language,
            element_type="structured_text",
            section_hint="",
        )
